#!/usr/bin/env python3
"""
BANA255 Final Presentation Generator
Creates a 10-minute professional PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import os

# Presentation settings
SLIDE_WIDTH = Inches(13.333)  # 16:9 widescreen
SLIDE_HEIGHT = Inches(7.5)

# RIT Orange color scheme
RIT_ORANGE = RGBColor(0xF7, 0x6C, 0x00)  # #F76C00
RIT_DARK = RGBColor(0x2D, 0x2D, 0x2D)     # Dark gray
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)

def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Orange header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RIT_ORANGE
    shape.line.fill.background()

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "THE BEST PIZZA NEAR RIT"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.8))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "A Data-Driven Market Entry Strategy"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Team info
    team = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12), Inches(1))
    tf = team.text_frame
    p = tf.paragraphs[0]
    p.text = "BANA255 Business Analytics | Fall 2025"
    p.font.size = Pt(24)
    p.font.color.rgb = RIT_DARK
    p.alignment = PP_ALIGN.CENTER

    # Key stat preview
    stat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(4.5), Inches(6), Inches(2))
    stat_box.fill.solid()
    stat_box.fill.fore_color.rgb = LIGHT_GRAY
    stat_box.line.color.rgb = RIT_ORANGE
    stat_box.line.width = Pt(3)

    stat = slide.shapes.add_textbox(Inches(3.5), Inches(4.7), Inches(6), Inches(1.8))
    tf = stat.text_frame
    p = tf.paragraphs[0]
    p.text = "161 Students Surveyed"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RIT_ORANGE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "52 Visualizations | 3-Component Model"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RIT_DARK
    p2.alignment = PP_ALIGN.CENTER

    return slide

def add_section_header(prs, title_text, subtitle_text=""):
    """Section header slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Left orange bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RIT_ORANGE
    shape.line.fill.background()

    # Title
    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RIT_DARK

    if subtitle_text:
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(11), Inches(1))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(24)
        p.font.color.rgb = RIT_ORANGE

    return slide

def add_content_slide(prs, title_text, bullets, image_path=None):
    """Standard content slide with optional image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RIT_ORANGE
    shape.line.fill.background()

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content positioning
    if image_path and os.path.exists(image_path):
        # Split layout: text left, image right
        content_left = Inches(0.5)
        content_width = Inches(5.5)
        img_left = Inches(6.5)
        img_width = Inches(6.3)
        img_height = Inches(5.5)

        # Add image
        slide.shapes.add_picture(image_path, img_left, Inches(1.5), width=img_width)
    else:
        content_left = Inches(0.5)
        content_width = Inches(12)

    # Bullets
    if bullets:
        content = slide.shapes.add_textbox(content_left, Inches(1.5), content_width, Inches(5.5))
        tf = content.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(22)
            p.font.color.rgb = RIT_DARK
            p.space_after = Pt(12)

    return slide

def add_two_column_slide(prs, title_text, left_title, left_bullets, right_title, right_bullets):
    """Two column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RIT_ORANGE
    shape.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column title
    lt = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.6))
    tf = lt.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RIT_ORANGE

    # Left column content
    lc = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.5))
    tf = lc.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RIT_DARK
        p.space_after = Pt(8)

    # Right column title
    rt = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.8), Inches(0.6))
    tf = rt.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RIT_ORANGE

    # Right column content
    rc = slide.shapes.add_textbox(Inches(7), Inches(2.2), Inches(5.8), Inches(4.5))
    tf = rc.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RIT_DARK
        p.space_after = Pt(8)

    return slide

def add_big_stat_slide(prs, stat_value, stat_label, subtitle=""):
    """Slide with big statistic"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background accent
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(1.5), Inches(7), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RIT_ORANGE
    shape.line.fill.background()

    # Big number
    stat = slide.shapes.add_textbox(Inches(3), Inches(1.8), Inches(7), Inches(2))
    tf = stat.text_frame
    p = tf.paragraphs[0]
    p.text = stat_value
    p.font.size = Pt(96)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Label
    label = slide.shapes.add_textbox(Inches(3), Inches(4), Inches(7), Inches(1))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = stat_label
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11), Inches(0.8))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = RIT_DARK
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_image_slide(prs, title_text, image_path, caption=""):
    """Full image slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Smaller header
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RIT_ORANGE
    shape.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Image - centered and large
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1.5), Inches(1.1), width=Inches(10))

    if caption:
        cap = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
        tf = cap.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = RIT_DARK
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_winner_slide(prs):
    """The Winner announcement slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full orange background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RIT_ORANGE
    shape.line.fill.background()

    # "THE WINNER IS..."
    pre = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(1))
    tf = pre.text_frame
    p = tf.paragraphs[0]
    p.text = "THE WINNER IS..."
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Winner name
    winner = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12), Inches(1.5))
    tf = winner.text_frame
    p = tf.paragraphs[0]
    p.text = "DOMINO'S PIZZA"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Stats box
    stats_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(4), Inches(8), Inches(2.5))
    stats_box.fill.solid()
    stats_box.fill.fore_color.rgb = WHITE
    stats_box.line.fill.background()

    stats = slide.shapes.add_textbox(Inches(2.5), Inches(4.2), Inches(8), Inches(2.3))
    tf = stats.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "27.3% Market Share (44 votes)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RIT_ORANGE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Threat Score: 73.9/100"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RIT_DARK
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Statistical Significance: p = 0.0001"
    p3.font.size = Pt(20)
    p3.font.color.rgb = RIT_DARK
    p3.alignment = PP_ALIGN.CENTER

    return slide

def add_swot_slide(prs):
    """SWOT Analysis slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RIT_ORANGE
    shape.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "SWOT ANALYSIS: Market Entry Opportunity"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Quadrants
    quad_width = Inches(5.8)
    quad_height = Inches(2.8)

    # Colors for quadrants
    colors = [
        RGBColor(0x4C, 0xAF, 0x50),  # Green - Strengths
        RGBColor(0xFF, 0x98, 0x00),  # Amber - Weaknesses
        RGBColor(0x21, 0x96, 0xF3),  # Blue - Opportunities
        RGBColor(0xF4, 0x43, 0x36),  # Red - Threats
    ]

    titles = ["STRENGTHS", "WEAKNESSES", "OPPORTUNITIES", "THREATS"]
    contents = [
        ["84% prefer local quality", "Fresh, authentic positioning", "Untapped 'persuadable' segment"],
        ["New entrant (no brand recognition)", "Higher operating costs", "Limited delivery infrastructure"],
        ["46pp preference-action gap", "74 students ready to switch", "$9.38 avg side order spend"],
        ["Domino's convenience dominance", "Chain price competition", "Student budget constraints"]
    ]

    positions = [
        (Inches(0.5), Inches(1.2)),
        (Inches(7), Inches(1.2)),
        (Inches(0.5), Inches(4.2)),
        (Inches(7), Inches(4.2))
    ]

    for i, (pos, color, title_t, content) in enumerate(zip(positions, colors, titles, contents)):
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos[0], pos[1], quad_width, quad_height)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = color
        box.line.width = Pt(3)

        # Title
        t = slide.shapes.add_textbox(pos[0] + Inches(0.2), pos[1] + Inches(0.1), quad_width - Inches(0.4), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = title_t
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color

        # Content
        c = slide.shapes.add_textbox(pos[0] + Inches(0.2), pos[1] + Inches(0.6), quad_width - Inches(0.4), Inches(2))
        tf = c.text_frame
        tf.word_wrap = True
        for j, item in enumerate(content):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(14)
            p.font.color.rgb = RIT_DARK

    return slide

def add_recommendation_slide(prs):
    """Final recommendation slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RIT_ORANGE
    shape.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "STRATEGIC RECOMMENDATION"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Main positioning statement
    pos_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(11), Inches(1.3))
    pos_box.fill.solid()
    pos_box.fill.fore_color.rgb = RIT_ORANGE
    pos_box.line.fill.background()

    pos = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(11), Inches(1))
    tf = pos.text_frame
    p = tf.paragraphs[0]
    p.text = '"LOCAL QUALITY AT CHAIN CONVENIENCE"'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Four pillars
    pillars = [
        ("PRODUCT", "Exceptional taste (94% priority)\nFresh ingredients, balanced flavor"),
        ("PRICE", "$18-20 for 16\" pizza\nOptimal Van Westendorp point"),
        ("SERVICE", "Fast pickup (<22 min)\nReliable delivery times"),
        ("TARGET", "Students with transportation\n2.7x more likely to choose local")
    ]

    pillar_width = Inches(2.8)
    start_x = Inches(0.5)

    for i, (title_t, content) in enumerate(pillars):
        x = start_x + i * (pillar_width + Inches(0.3))

        # Pillar box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.1), pillar_width, Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = RIT_ORANGE
        box.line.width = Pt(2)

        # Title
        t = slide.shapes.add_textbox(x, Inches(3.3), pillar_width, Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = title_t
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RIT_ORANGE
        p.alignment = PP_ALIGN.CENTER

        # Content
        c = slide.shapes.add_textbox(x + Inches(0.15), Inches(3.9), pillar_width - Inches(0.3), Inches(2.5))
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(14)
        p.font.color.rgb = RIT_DARK
        p.alignment = PP_ALIGN.CENTER

    # Bottom call to action
    cta = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
    tf = cta.text_frame
    p = tf.paragraphs[0]
    p.text = "Capture the 74 'persuadable' students ready to switch from chains to quality local pizza"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RIT_DARK
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_closing_slide(prs):
    """Final Q&A slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full orange
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RIT_ORANGE
    shape.line.fill.background()

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "QUESTIONS?"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Summary
    summary = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(2))
    tf = summary.text_frame
    p = tf.paragraphs[0]
    p.text = "The Best Pizza = Taste + Value + Speed"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Winner: Domino's | Opportunity: 74 Persuadable Students"
    p2.font.size = Pt(24)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    return slide


def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    output_dir = "outputs"

    # ==================== SLIDES ====================

    # SLIDE 1: Title
    add_title_slide(prs)
    print("Created: Title Slide")

    # SLIDE 2: ACTION Overview Section
    add_section_header(prs, "THE ACTION FRAMEWORK", "Our Analytical Approach")
    print("Created: ACTION Section Header")

    # SLIDE 3: ACTION Overview Content
    add_two_column_slide(prs,
        "ACTION Framework Application",
        "ASK & CLARIFY",
        [
            "Question: What defines 'best pizza' for RIT students?",
            "161 consented survey responses",
            "98.6% data completeness, 92/100 quality score",
            "9 pizza quality factors analyzed"
        ],
        "THINK & INFORM",
        [
            "Weighted importance model (Friedman test)",
            "ML prediction: RF, GBM, Logistic (71% acc)",
            "Competitive ranking with composite scores",
            "Van Westendorp price sensitivity analysis"
        ]
    )
    print("Created: ACTION Overview")

    # SLIDE 4: Model Framework Section
    add_section_header(prs, "OUR 'BEST PIZZA' MODEL", "Three Components, One Answer")
    print("Created: Model Section Header")

    # SLIDE 5: Component 1 - Importance Weights
    add_content_slide(prs,
        "Component 1: What Makes Pizza 'Best'?",
        [
            "Taste & Flavor: 14.1% weight (94% rate highly important)",
            "Balance: 12.6% | Crust: 12.2% | Freshness: 12.1%",
            "Price: 12.0% (critical value factor)",
            "Core Quality factors = 51% of total weight",
            "Friedman Test: X² = 435.85, p < 0.001"
        ],
        f"{output_dir}/fig7_importance_weights.png"
    )
    print("Created: Importance Weights Slide")

    # SLIDE 6: Component 2 - Behavioral Model
    add_content_slide(prs,
        "Component 2: Predicting Local vs Chain Choice",
        [
            "71.1% accuracy WITHOUT circular features",
            "Excluded 'states prefer local' (tautological)",
            "Top predictors: Pickup time, Max price, Transportation",
            "Students with cars: 2.7x more likely to choose local",
            "Key insight: Convenience drives actual behavior"
        ],
        f"{output_dir}/fig13_feature_importance.png"
    )
    print("Created: Behavioral Model Slide")

    # SLIDE 7: Component 3 - Competitive Ranking
    add_image_slide(prs,
        "Component 3: Competitive Threat Ranking",
        f"{output_dir}/fig8_competitive_ranking.png",
        "Composite score: Market Share (30%) + Loyalty (25%) + Profile Match (25%) + Local Capture (20%)"
    )
    print("Created: Competitive Ranking Slide")

    # SLIDE 8: THE WINNER
    add_winner_slide(prs)
    print("Created: Winner Announcement Slide")

    # SLIDE 9: The Local-Chain Paradox
    add_big_stat_slide(prs,
        "46pp",
        "THE LOCAL-CHAIN PARADOX",
        "84% SAY they prefer local pizza, but only 38% CHOOSE local pizza"
    )
    print("Created: Paradox Stat Slide")

    # SLIDE 10: Paradox Visual
    add_content_slide(prs,
        "The Preference-Action Gap Explained",
        [
            "74 students WANT local but choose chains",
            "Why? Convenience constraints:",
            "  - Transportation barriers (no car)",
            "  - Time constraints (chains are faster)",
            "  - Price perception (chains seem cheaper)",
            "This gap = Our market opportunity"
        ],
        f"{output_dir}/fig9_local_chain_paradox.png"
    )
    print("Created: Paradox Explanation Slide")

    # SLIDE 11: Strategic Section Header
    add_section_header(prs, "MARKET ENTRY STRATEGY", "Capturing the Opportunity")
    print("Created: Strategy Section Header")

    # SLIDE 12: SWOT Analysis
    add_swot_slide(prs)
    print("Created: SWOT Analysis Slide")

    # SLIDE 13: Price Strategy
    add_content_slide(prs,
        "Optimal Pricing: Van Westendorp Analysis",
        [
            "Optimal Price Point: $18-20 for 16\" pizza",
            "Acceptable range: $15-22",
            "Below $15: Perceived as low quality",
            "Above $22: Too expensive for students",
            "Side order opportunity: $9.38 average spend"
        ],
        f"{output_dir}/fig35_van_westendorp.png"
    )
    print("Created: Pricing Slide")

    # SLIDE 14: Target Segments
    add_content_slide(prs,
        "Target Customer Segments",
        [
            "Primary: Quality Seekers (29%)",
            "  - Prioritize taste over price",
            "  - Higher price tolerance ($20+)",
            "Secondary: Value Hunters (35%)",
            "  - Price-sensitive but quality-aware",
            "  - Target with combo deals",
            "Key: Students with transportation access"
        ],
        f"{output_dir}/fig16_customer_profiles.png"
    )
    print("Created: Segments Slide")

    # SLIDE 15: Side Order Revenue
    add_content_slide(prs,
        "Side Order Revenue Opportunity",
        [
            "Garlic knots: 65% would likely order",
            "Wings: 53% would likely order",
            "Breadsticks: 41% would likely order",
            "Average side spend: $9.38 per order",
            "Strategic recommendation: Bundle with pizza deals"
        ],
        f"{output_dir}/fig47_side_order_rankings.png"
    )
    print("Created: Side Orders Slide")

    # SLIDE 16: Final Recommendation
    add_recommendation_slide(prs)
    print("Created: Recommendation Slide")

    # SLIDE 17: Key Metrics Summary
    add_two_column_slide(prs,
        "Key Metrics Summary",
        "MARKET INSIGHTS",
        [
            "161 students surveyed (98.6% complete)",
            "Domino's leads: 27.3% market share",
            "84% prefer local, 38% choose local",
            "46 percentage point gap = opportunity",
            "71.1% prediction accuracy"
        ],
        "SUCCESS FACTORS",
        [
            "Taste is #1 (14.1% weight, 94% priority)",
            "Speed matters: <22 min pickup expected",
            "Price sweet spot: $18-20",
            "Side orders: +$9.38 per ticket",
            "Target: Students with transportation"
        ]
    )
    print("Created: Metrics Summary Slide")

    # SLIDE 18: Closing
    add_closing_slide(prs)
    print("Created: Closing Slide")

    # Save presentation
    output_path = "outputs/BANA255_Best_Pizza_Presentation.pptx"
    prs.save(output_path)
    print(f"\n✓ Presentation saved to: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")

    return output_path

if __name__ == "__main__":
    main()
